from pptx import Presentation

prs = Presentation('Django –Module 3.pptx')
slides = prs.slides

print(f"Total slides: {len(slides)}\n")

# Show all slides from 26 onwards
for slide_num in range(25, len(slides)):
    slide = slides[slide_num]
    print(f"\nSlide {slide_num + 1}:")
    print("=" * 70)
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text)
    print("=" * 70)


